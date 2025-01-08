import os
import pytesseract
import cv2
from pptx import Presentation
from pptx.util import Inches

def ocr_extract_text_pysseract(image_path):
    print("Extracting Text from Image:", image_path)
    # Load the image from the given file path
    image = cv2.imread(image_path)
    # Convert the image to grayscale
    gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

    # Extract text from the image
    text = pytesseract.image_to_string(gray_image)
    return text


def add_extracted_text_to_powerpoint(pptx_path,template_path, text):
    # Check if the PowerPoint file exists
    if os.path.exists(pptx_path):
        # Load the existing PowerPoint presentation
        prs = Presentation(pptx_path)
    else:
        # Create a new PowerPoint presentation
        prs = Presentation(template_path)

    # Add a new slide to the presentation
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add the extracted text to the slide
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.text = text

    # Save the updated presentation
    prs.save(pptx_path)
def numerical_sort(value):
    import re
    numbers = re.findall(r'\d+', value)
    return [int(num) for num in numbers]

def main():
    ####### Configuration #######
    template_path = "sample.pptx" # PowerPoint template
    extract_path = "image1" # Folder containing the images

    # Output folder path
    pptx_path = f"/output/{extract_path}_extract.pptx"
    #############################

    image_folder_path = f"input/{extract_path}"
    # Numerically sort the filenames
    # YES, this is a hacky way to sort the filenames - CHANGE If you have a better way
    filenames = sorted(os.listdir(image_folder_path), key=numerical_sort)
    # Iterate through the image folder
    for filename in filenames:
        if filename.endswith(".jpeg") or filename.endswith(".jpg"):
            print("Processing File:", filename)
            # Construct the full path to the image file
            image_path = os.path.join(image_folder_path, filename)
            # Extract text from the image
            text = ocr_extract_text_pysseract(image_path)
            # Add the extracted text to the PowerPoint presentation
            add_extracted_text_to_powerpoint(pptx_path,template_path, text)
    print("Presentation Updated")  # Debugging statement

if __name__ == "__main__":
    main()

